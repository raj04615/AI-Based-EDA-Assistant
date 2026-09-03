"""
Ingestion Engine — Multi-format document parsing (PDF, DOCX, PPTX, XLSX/CSV, TXT/MD, Image OCR),
sentence-aware chunking, and background embedding/indexing pipeline.
"""

import re
import os
import zipfile
import traceback
from datetime import datetime, timezone
from typing import List, Dict, Any

from pypdf import PdfReader

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    pdfplumber = None
    HAS_PDFPLUMBER = False

import config
import database

# Lazy-loaded embedding model handle
_embedding_model = None

def get_embedding_model():
    """Load SentenceTransformer model lazily."""
    global _embedding_model
    if _embedding_model is None:
        from sentence_transformers import SentenceTransformer
        _embedding_model = SentenceTransformer(config.EMBEDDING_MODEL)
    return _embedding_model


def detect_file_type(file_path: str) -> str:
    """
    Detect file format using magic bytes signature inspection with fallback to file extension.
    Returns one of: 'pdf', 'docx', 'pptx', 'excel', 'csv', 'txt', 'md', 'image', 'zip', or 'unknown'.
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    # Fast reject unsupported extensions (e.g. .exe, .bin, .dll)
    valid_exts = {".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv", ".txt", ".log", ".md", ".markdown", ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".zip"}
    if ext not in valid_exts:
        return "unknown"
    
    try:
        with open(file_path, "rb") as f:
            header = f.read(2048)
    except Exception:
        header = b""

    # Magic Bytes Inspections
    if header.startswith(b"%PDF-"):
        return "pdf"

    if header.startswith(b"\x89PNG\r\n\x1a\n") or header.startswith(b"\xff\xd8\xff") or header.startswith(b"GIF8") or header.startswith(b"BM"):
        return "image"

    # Zip-based archives (ZIP, DOCX, PPTX, XLSX)
    if header.startswith(b"PK\x03\x04"):
        if ext == ".docx":
            return "docx"
        elif ext == ".pptx":
            return "pptx"
        elif ext in (".xlsx", ".xls"):
            return "excel"
        elif ext == ".zip":
            return "zip"
        else:
            try:
                with zipfile.ZipFile(file_path, 'r') as z:
                    names = z.namelist()
                    if "word/document.xml" in names:
                        return "docx"
                    if "ppt/presentation.xml" in names:
                        return "pptx"
                    if "xl/workbook.xml" in names:
                        return "excel"
            except Exception:
                pass
            return "zip"

    # Extension Fallbacks
    if ext == ".pdf":
        return "pdf"
    elif ext == ".docx":
        return "docx"
    elif ext == ".pptx":
        return "pptx"
    elif ext in (".xlsx", ".xls"):
        return "excel"
    elif ext == ".csv":
        return "csv"
    elif ext in (".txt", ".log"):
        return "txt"
    elif ext in (".md", ".markdown"):
        return "md"
    elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff"):
        return "image"
    elif ext == ".zip":
        return "zip"

    # Plain text detection fallback
    try:
        header.decode('utf-8')
        return "md" if ext == ".md" else "txt"
    except UnicodeDecodeError:
        pass

    return "unknown"


def extract_table_as_markdown(table: List[List[Any]]) -> str:
    """Convert a raw table matrix into a Markdown formatted table."""
    if not table or not any(table):
        return ""
    
    clean_rows = []
    for row in table:
        if row and any(cell is not None and str(cell).strip() for cell in row):
            clean_rows.append([str(cell).strip().replace("\n", " ") if cell is not None else "" for cell in row])
            
    if not clean_rows:
        return ""
        
    headers = clean_rows[0]
    markdown = "| " + " | ".join(headers) + " |\n"
    markdown += "| " + " | ".join(["---"] * len(headers)) + " |\n"
    
    for row in clean_rows[1:]:
        if len(row) < len(headers):
            row.extend([""] * (len(headers) - len(row)))
        elif len(row) > len(headers):
            row = row[:len(headers)]
        markdown += "| " + " | ".join(row) + " |\n"
        
    return markdown


# ── Format-Specific Extractors ────────────────────────────────────────

def extract_pdf(pdf_path: str) -> List[Dict[str, Any]]:
    """Extract text and structured tables per page from PDF files."""
    pages_data = []
    tables_per_page = {}
    
    if HAS_PDFPLUMBER and pdfplumber:
        try:
            with pdfplumber.open(pdf_path) as pdf:
                for i, page in enumerate(pdf.pages, start=1):
                    extracted_tables = page.extract_tables()
                    md_tables = []
                    for tbl in extracted_tables:
                        md_tbl = extract_table_as_markdown(tbl)
                        if md_tbl:
                            md_tables.append(md_tbl)
                    if md_tables:
                        tables_per_page[i] = md_tables
        except Exception as e:
            print(f"[INGESTION] pdfplumber table extraction warning: {e}")

    reader = PdfReader(pdf_path)
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        text = text.strip()
        tables = tables_per_page.get(i, [])
        if text or tables:
            pages_data.append({
                "page": i,
                "page_label": f"Page {i}",
                "text": text,
                "tables": tables
            })
            
    return pages_data


def extract_docx(docx_path: str) -> List[Dict[str, Any]]:
    """Extract paragraphs and tables from Microsoft Word (.docx) files."""
    import docx
    doc = docx.Document(docx_path)
    
    sections_data = []
    current_text = []
    current_tables = []
    section_index = 1

    for elem in doc.element.body:
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        if tag == "p":
            para = docx.text.paragraph.Paragraph(elem, doc)
            text = para.text.strip()
            if text:
                current_text.append(text)
                if sum(len(t) for t in current_text) >= 1500:
                    sections_data.append({
                        "page": section_index,
                        "page_label": f"Section {section_index}",
                        "text": "\n\n".join(current_text),
                        "tables": current_tables
                    })
                    section_index += 1
                    current_text = []
                    current_tables = []
        elif tag == "tbl":
            tbl = docx.table.Table(elem, doc)
            matrix = []
            for row in tbl.rows:
                matrix.append([cell.text.strip() for cell in row.cells])
            md_tbl = extract_table_as_markdown(matrix)
            if md_tbl:
                current_tables.append(md_tbl)

    if current_text or current_tables:
        sections_data.append({
            "page": section_index,
            "page_label": f"Section {section_index}",
            "text": "\n\n".join(current_text),
            "tables": current_tables
        })

    return sections_data


def extract_pptx(pptx_path: str) -> List[Dict[str, Any]]:
    """Extract slide text, tables, and speaker notes from PowerPoint (.pptx) files."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    
    slides_data = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []
        slide_tables = []
        
        # Extract text & tables from shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_text_parts.append(text)
            elif shape.has_table:
                matrix = []
                for row in shape.table.rows:
                    matrix.append([cell.text.strip() for cell in row.cells])
                md_tbl = extract_table_as_markdown(matrix)
                if md_tbl:
                    slide_tables.append(md_tbl)

        # Extract speaker notes if available
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text_parts.append(f"[Speaker Notes]: {notes}")

        combined_text = "\n\n".join(slide_text_parts)
        if combined_text or slide_tables:
            slides_data.append({
                "page": i,
                "page_label": f"Slide {i}",
                "text": combined_text,
                "tables": slide_tables
            })

    return slides_data


def extract_excel_csv(file_path: str, file_type: str) -> List[Dict[str, Any]]:
    """Extract sheets and dataframes from Excel (.xlsx) or CSV files into Markdown tables."""
    import pandas as pd
    pages_data = []

    if file_type == "csv":
        df = None
        for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1"):
            try:
                df = pd.read_csv(file_path, encoding=enc)
                break
            except Exception:
                continue
        if df is None:
            df = pd.read_csv(file_path, encoding_errors="replace")
        sheets = {"Data": df}
    else:
        sheets = pd.read_excel(file_path, sheet_name=None)

    sec_idx = 1
    for sheet_name, df in sheets.items():
        if df.empty:
            continue
        
        # Fill NaN values with empty string
        df_clean = df.fillna("")
        
        # Chunk large dataframes into 50-row blocks
        batch_size = 50
        num_rows = len(df_clean)
        
        for start_row in range(0, num_rows, batch_size):
            end_row = min(start_row + batch_size, num_rows)
            sub_df = df_clean.iloc[start_row:end_row]
            
            headers = [str(col) for col in sub_df.columns]
            matrix = [headers] + sub_df.astype(str).values.tolist()
            table_md = extract_table_as_markdown(matrix)
            
            if table_md:
                row_label = f"Sheet '{sheet_name}' (Rows {start_row + 1}-{end_row})"
                pages_data.append({
                    "page": sec_idx,
                    "page_label": row_label,
                    "text": f"Data excerpt from {row_label}:",
                    "tables": [table_md]
                })
                sec_idx += 1

    return pages_data


def extract_txt_md(txt_path: str) -> List[Dict[str, Any]]:
    """Read plain text or Markdown files with encoding fallbacks."""
    content = ""
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1"):
        try:
            with open(txt_path, "r", encoding=enc) as f:
                content = f.read()
            break
        except Exception:
            continue

    if not content:
        try:
            with open(txt_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
        except Exception:
            pass

    if not content.strip():
        return []

    # Split long text by double newlines into ~1500 char sections
    paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
    pages_data = []
    current_paras = []
    sec_idx = 1

    for p in paragraphs:
        current_paras.append(p)
        if sum(len(x) for x in current_paras) >= 1500:
            pages_data.append({
                "page": sec_idx,
                "page_label": f"Page {sec_idx}",
                "text": "\n\n".join(current_paras),
                "tables": []
            })
            sec_idx += 1
            current_paras = []

    if current_paras:
        pages_data.append({
            "page": sec_idx,
            "page_label": f"Page {sec_idx}",
            "text": "\n\n".join(current_paras),
            "tables": []
        })

    return pages_data


def extract_image(img_path: str) -> List[Dict[str, Any]]:
    """Extract text from image scans or screenshots using pytesseract OCR."""
    from PIL import Image
    import pytesseract

    try:
        img = Image.open(img_path)
        ocr_text = pytesseract.image_to_string(img).strip()
    except Exception as e:
        if "tesseract is not installed" in str(e).lower() or "not in your PATH" in str(e):
            raise RuntimeError(
                "Tesseract OCR is not installed on the system environment. "
                "Please install Tesseract OCR binary or use text-based documents."
            ) from e
        raise

    if not ocr_text:
        return []

    return [{
        "page": 1,
        "page_label": "Image 1",
        "text": ocr_text,
        "tables": []
    }]


def extract_by_type(file_path: str, file_type: str) -> List[Dict[str, Any]]:
    """Central router calling the dedicated extractor based on detected format."""
    if file_type == "pdf":
        return extract_pdf(file_path)
    elif file_type == "docx":
        return extract_docx(file_path)
    elif file_type == "pptx":
        return extract_pptx(file_path)
    elif file_type in ("excel", "csv"):
        return extract_excel_csv(file_path, file_type)
    elif file_type in ("txt", "md"):
        return extract_txt_md(file_path)
    elif file_type == "image":
        return extract_image(file_path)
    else:
        raise ValueError(f"Unsupported document format type '{file_type}'.")


# ── Markdown Normalization Engine ─────────────────────────────────────

def _preserve_list_formatting(text: str) -> str:
    """Ensure bulleted and numbered list items retain clean linebreaks."""
    lines = text.split("\n")
    formatted_lines = []
    in_list = False

    for line in lines:
        stripped = line.strip()
        is_list_item = bool(re.match(r'^([\*\-\+]\s+|\d+[\.\)]\s+)', stripped))
        if is_list_item:
            if not in_list and formatted_lines and formatted_lines[-1] != "":
                formatted_lines.append("")
            formatted_lines.append(stripped)
            in_list = True
        else:
            if in_list and stripped != "":
                in_list = False
            formatted_lines.append(line)

    return "\n".join(formatted_lines)


def normalize_to_markdown(
    extracted_content: List[Dict[str, Any]],
    file_type: str,
    filename: str,
    doc_id: str
) -> str:
    """
    Convert raw extracted content (text blocks, tables, slide/page metadata) into a single,
    canonical well-structured Markdown (.md) document with YAML front matter header.
    """
    now_iso = datetime.now(timezone.utc).isoformat()

    # Determine count metric (page_count / slide_count / row_count)
    count_key = "page_count"
    count_val = len(extracted_content)

    if file_type == "pptx":
        count_key = "slide_count"
        count_val = len(extracted_content)
    elif file_type in ("excel", "csv"):
        count_key = "row_count"
        total_rows = 0
        for item in extracted_content:
            for tbl in item.get("tables", []):
                lines = [l for l in tbl.splitlines() if l.strip() and not l.strip().startswith("| ---")]
                if lines:
                    total_rows += max(0, len(lines) - 1)
        count_val = total_rows if total_rows > 0 else len(extracted_content)

    front_matter = (
        "---\n"
        f"title: {filename}\n"
        f"source_filename: {filename}\n"
        f"file_type: {file_type}\n"
        f"doc_id: {doc_id}\n"
        f"ingested_at: {now_iso}\n"
        f"{count_key}: {count_val}\n"
        "---\n\n"
    )

    body_parts = []

    for item in extracted_content:
        page_num = item.get("page", 1)
        page_label = item.get("page_label", f"Page {page_num}")
        text = item.get("text", "").strip()
        tables = item.get("tables", [])

        # Format section headings
        if file_type == "image":
            section_header = "## Extracted Text (OCR)"
        elif file_type in ("excel", "csv"):
            section_header = f"## {page_label}"
        elif file_type == "pptx":
            section_header = f"## {page_label}"
        else:
            # Check if text already starts with a Markdown header
            if text and (text.startswith("# ") or text.startswith("## ") or text.startswith("### ")):
                section_header = ""
            else:
                section_header = f"## {page_label}"

        section_block = []
        if section_header:
            section_block.append(section_header)

        if text:
            formatted_text = _preserve_list_formatting(text)
            section_block.append(formatted_text)

        if tables:
            for tbl in tables:
                if tbl.strip():
                    section_block.append(tbl.strip())

        if section_block:
            body_parts.append("\n\n".join(section_block))

    return front_matter + "\n\n".join(body_parts) + "\n"


# ── Canonical Markdown Chunker ───────────────────────────────────────

def chunk_markdown(
    markdown_content: str,
    chunk_size: int = config.CHUNK_SIZE,
    overlap: int = config.CHUNK_OVERLAP
) -> List[Dict[str, Any]]:
    """
    Standardized sentence and table-aware chunker operating on canonical Markdown content.
    Preserves table structure and structural headings for consistent chunking across formats.
    """
    content = markdown_content
    # Strip YAML front matter if present
    if content.startswith("---"):
        parts = content.split("---", 2)
        if len(parts) >= 3:
            content = parts[2].strip()

    section_pattern = re.compile(r'^(#+\s+.*)$', re.MULTILINE)
    splits = section_pattern.split(content)

    chunks = []
    current_label = "Document Content"
    current_page = 1

    idx = 0
    if splits and not splits[0].startswith("#"):
        initial_block = splits[0].strip()
        idx = 1
    else:
        initial_block = ""

    blocks_to_process = []
    if initial_block:
        blocks_to_process.append((current_label, current_page, initial_block))

    while idx < len(splits):
        heading = splits[idx].strip()
        body = splits[idx + 1].strip() if idx + 1 < len(splits) else ""
        idx += 2

        clean_heading = heading.lstrip("#").strip()
        current_label = clean_heading if clean_heading else "Section"

        page_match = re.search(r'\b(?:Page|Slide|Section)\s+(\d+)\b', clean_heading, re.IGNORECASE)
        if page_match:
            current_page = int(page_match.group(1))

        if body:
            blocks_to_process.append((current_label, current_page, body))

    for label, page_num, body_text in blocks_to_process:
        lines = body_text.split("\n")
        para_lines = []
        tbl_lines = []
        in_table = False

        def flush_table(table_rows):
            if not table_rows:
                return
            tbl_md = "\n".join(table_rows)
            chunks.append({
                "text": f"Table Excerpt [{label}]:\n{tbl_md}",
                "pages": [page_num],
                "page_labels": [label],
                "is_table": True
            })

        def flush_paragraphs(p_lines):
            if not p_lines:
                return
            raw_p_text = "\n".join(p_lines).strip()
            if not raw_p_text:
                return

            paragraphs = [p.strip() for p in raw_p_text.split("\n\n") if p.strip()]
            current_chunk = ""

            for para in paragraphs:
                if len(current_chunk) + len(para) + 2 <= chunk_size:
                    current_chunk = f"{current_chunk}\n\n{para}".strip()
                else:
                    if current_chunk:
                        chunks.append({
                            "text": current_chunk,
                            "pages": [page_num],
                            "page_labels": [label],
                            "is_table": False
                        })

                    if len(para) > chunk_size:
                        sentences = re.split(r'(?<=[.!?])\s+', para)
                        sub_chunk = ""
                        for sentence in sentences:
                            if len(sub_chunk) + len(sentence) + 1 <= chunk_size:
                                sub_chunk = f"{sub_chunk} {sentence}".strip()
                            else:
                                if sub_chunk:
                                    chunks.append({
                                        "text": sub_chunk,
                                        "pages": [page_num],
                                        "page_labels": [label],
                                        "is_table": False
                                    })
                                sub_chunk = sentence
                        if sub_chunk:
                            current_chunk = sub_chunk
                    else:
                        current_chunk = para

            if current_chunk:
                chunks.append({
                    "text": current_chunk,
                    "pages": [page_num],
                    "page_labels": [label],
                    "is_table": False
                })

        for line in lines:
            if line.strip().startswith("|"):
                if not in_table:
                    flush_paragraphs(para_lines)
                    para_lines = []
                    in_table = True
                tbl_lines.append(line)
            else:
                if in_table:
                    flush_table(tbl_lines)
                    tbl_lines = []
                    in_table = False
                para_lines.append(line)

        if in_table:
            flush_table(tbl_lines)
        else:
            flush_paragraphs(para_lines)

    return chunks


def chunk_text_sentence_aware(
    pages_data: List[Dict[str, Any]],
    chunk_size: int = config.CHUNK_SIZE,
    overlap: int = config.CHUNK_OVERLAP
) -> List[Dict[str, Any]]:
    """
    Sentence and paragraph-aware chunker operating on pages_data.
    Maintained for backward compatibility and unit tests.
    """
    chunks = []
    
    for page_data in pages_data:
        page_num = page_data["page"]
        page_label = page_data.get("page_label", f"Page {page_num}")
        
        # 1. Process Table Chunks
        for tbl_idx, table_md in enumerate(page_data.get("tables", [])):
            chunks.append({
                "text": f"Table Excerpt [{page_label}]:\n{table_md}",
                "pages": [page_num],
                "page_labels": [page_label],
                "is_table": True
            })

        # 2. Process Narrative Text Chunks
        raw_text = page_data.get("text", "")
        if not raw_text:
            continue
            
        paragraphs = [p.strip() for p in raw_text.split("\n\n") if p.strip()]
        current_chunk = ""
        
        for para in paragraphs:
            if len(current_chunk) + len(para) + 2 <= chunk_size:
                current_chunk = f"{current_chunk}\n\n{para}".strip()
            else:
                if current_chunk:
                    chunks.append({
                        "text": current_chunk,
                        "pages": [page_num],
                        "page_labels": [page_label],
                        "is_table": False
                    })
                    
                if len(para) > chunk_size:
                    sentences = re.split(r'(?<=[.!?])\s+', para)
                    sub_chunk = ""
                    for sentence in sentences:
                        if len(sub_chunk) + len(sentence) + 1 <= chunk_size:
                            sub_chunk = f"{sub_chunk} {sentence}".strip()
                        else:
                            if sub_chunk:
                                chunks.append({
                                    "text": sub_chunk,
                                    "pages": [page_num],
                                    "page_labels": [page_label],
                                    "is_table": False
                                })
                            sub_chunk = sentence
                    if sub_chunk:
                        current_chunk = sub_chunk
                else:
                    current_chunk = para
                    
        if current_chunk:
            chunks.append({
                "text": current_chunk,
                "pages": [page_num],
                "page_labels": [page_label],
                "is_table": False
            })

    return chunks


# ── Ingestion Background Pipeline ────────────────────────────────────

def run_ingestion_pipeline(doc_id: str, file_path: str, file_type: str = "pdf"):
    """
    Isolated async background task executing the full ingestion workflow:
    parse -> normalize to Markdown -> chunk canonical Markdown -> embed -> Pinecone upsert under namespace=doc_id.
    """
    try:
        # Phase 1: Parsing
        database.update_document_status(doc_id, "parsing")
        pages_data = extract_by_type(file_path, file_type)
        page_count = len(pages_data)
        database.update_document_status(doc_id, "parsing", page_count=page_count)

        if not pages_data:
            database.update_document_status(
                doc_id, "failed", error_message="No extractable text or tables found in document."
            )
            return

        # Phase 1.5: Markdown Normalization & Persistence
        filename = os.path.basename(file_path)
        markdown_str = normalize_to_markdown(pages_data, file_type, filename, doc_id)
        md_path = os.path.join(config.PROCESSED_DIR, f"{doc_id}.md")
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(markdown_str)

        database.update_document_status(doc_id, "parsing", markdown_path=md_path)

        # Phase 2: Chunking against canonical Markdown representation
        database.update_document_status(doc_id, "chunking")
        chunks = chunk_markdown(markdown_str)
        chunk_count = len(chunks)
        database.update_document_status(doc_id, "chunking", chunk_count=chunk_count)

        if not chunks:
            database.update_document_status(
                doc_id, "failed", error_message="Failed to generate chunks from document."
            )
            return

        # Phase 3: Embedding
        database.update_document_status(doc_id, "embedding")
        embedder = get_embedding_model()
        texts_to_embed = [c["text"] for c in chunks]
        embeddings = embedder.encode(texts_to_embed, show_progress_bar=False, normalize_embeddings=True).tolist()

        # Phase 4: Pinecone Vector Storage (Namespaced)
        database.update_document_status(doc_id, "indexing")
        
        from pinecone import Pinecone, ServerlessSpec
        pc = Pinecone(api_key=config.PINECONE_API_KEY)
        
        existing_indexes = {idx.name: idx for idx in pc.list_indexes()}
        if config.PINECONE_INDEX_NAME in existing_indexes:
            idx_info = existing_indexes[config.PINECONE_INDEX_NAME]
            if getattr(idx_info, "dimension", None) and idx_info.dimension != config.EMBEDDING_DIMENSION:
                print(f"[PINECONE] Dimension mismatch: index is {idx_info.dimension}, model requires {config.EMBEDDING_DIMENSION}. Recreating index...")
                pc.delete_index(config.PINECONE_INDEX_NAME)
                import time
                time.sleep(2)
                pc.create_index(
                    name=config.PINECONE_INDEX_NAME,
                    dimension=config.EMBEDDING_DIMENSION,
                    metric="cosine",
                    spec=ServerlessSpec(cloud="aws", region="us-east-1"),
                )
        else:
            pc.create_index(
                name=config.PINECONE_INDEX_NAME,
                dimension=config.EMBEDDING_DIMENSION,
                metric="cosine",
                spec=ServerlessSpec(cloud="aws", region="us-east-1"),
            )

        index = pc.Index(config.PINECONE_INDEX_NAME)
        
        vectors = []
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
            page_labels = chunk.get("page_labels", [str(p) for p in chunk["pages"]])
            vectors.append({
                "id": f"{doc_id}-chunk-{i}",
                "values": embedding,
                "metadata": {
                    "text": chunk["text"],
                    "pages": page_labels,
                    "is_table": chunk.get("is_table", False),
                    "doc_id": doc_id
                }
            })

        batch_size = 100
        for b_start in range(0, len(vectors), batch_size):
            batch = vectors[b_start : b_start + batch_size]
            index.upsert(vectors=batch, namespace=doc_id)

        database.update_document_status(doc_id, "ready")
        print(f"[INGESTION SUCCESS] Document {doc_id} ({file_type}) processed: {page_count} units, {chunk_count} chunks indexed.")

    except Exception as e:
        error_msg = f"Ingestion error: {str(e)}"
        print(f"[INGESTION ERROR] Document {doc_id} ({file_type}): {error_msg}")
        traceback.print_exc()
        database.update_document_status(doc_id, "failed", error_message=error_msg)
